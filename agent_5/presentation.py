import os
import io
import base64
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PresentationGenerator:
    def __init__(self):
        self.presentation = None
    
    def generate_presentation(self, title, content, creator_name=None, theme="modern"):
        """
        Generate a PowerPoint presentation based on content creator analysis
        
        Args:
            title (str): Presentation title
            content (dict): Content to include in the presentation
            creator_name (str, optional): Name of the content creator
            theme (str, optional): Theme for the presentation (modern, minimal, bold, etc.)
        
        Returns:
            bytes: PowerPoint presentation file as bytes
        """
        try:
            # Create a new PowerPoint presentation
            self.presentation = Presentation()
            
            # Set slide dimensions to widescreen (16:9)
            self.presentation.slide_width = Inches(13.33)
            self.presentation.slide_height = Inches(7.5)
            
            # Add title slide
            self._add_title_slide(title, creator_name)
            
            # Add overview slide
            self._add_overview_slide()
            
            # Add platform-specific slides based on available data
            if "youtube" in content and content["youtube"] and "error" not in content["youtube"]:
                self._add_youtube_slide(content["youtube"])
            
            if "tiktok" in content and content["tiktok"] and "error" not in content["tiktok"]:
                self._add_tiktok_slide(content["tiktok"])
            
            # Add suggestions slide
            if "suggestions" in content and content["suggestions"]:
                self._add_suggestions_slides(content["suggestions"])
            
            # Add next steps slide
            self._add_next_steps_slide()
            
            # Add thank you slide
            self._add_thank_you_slide()
            
            # Save presentation to a bytes buffer
            output = io.BytesIO()
            self.presentation.save(output)
            output.seek(0)
            
            return output.getvalue()
            
        except Exception as e:
            return {"error": f"Failed to generate presentation: {str(e)}"}
    
    def _add_title_slide(self, title, creator_name=None):
        """Add the title slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        if creator_name:
            subtitle.text = f"Performance Report for {creator_name}"
        else:
            subtitle.text = "Social Media Performance Report"
        
        # Add date
        current_date = datetime.now().strftime("%B %d, %Y")
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        date_text = date_box.text_frame
        date_text.text = f"Generated on {current_date}"
        
        return slide
    
    def _add_overview_slide(self):
        """Add an overview slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Overview"
        
        # Set content
        content_shape = slide.placeholders[1]
        content_text = content_shape.text_frame
        content_text.text = "This presentation provides an analysis of your social media performance across platforms and offers suggestions for optimization."
        
        p = content_text.add_paragraph()
        p.text = "The analysis includes:"
        
        p = content_text.add_paragraph()
        p.text = "• Performance metrics from your social media accounts"
        p.level = 1
        
        p = content_text.add_paragraph()
        p.text = "• Comparative analysis of platform performance"
        p.level = 1
        
        p = content_text.add_paragraph()
        p.text = "• Tailored optimization suggestions"
        p.level = 1
        
        p = content_text.add_paragraph()
        p.text = "• Recommended next steps for implementation"
        p.level = 1
        
        return slide
    
    def _add_youtube_slide(self, youtube_data):
        """Add YouTube statistics slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "YouTube Performance"
        
        # Clear default content placeholder
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx != 0:  # Skip title placeholder
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Create a table for statistics
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(12)
        height = Inches(1.2)
        
        # Add metric boxes
        self._add_metric_box(slide, "Subscribers", youtube_data.get("subscribers", "N/A"), Inches(0.5), Inches(2.0), Inches(3.8), Inches(1.5))
        self._add_metric_box(slide, "Total Views", youtube_data.get("total_views", "N/A"), Inches(4.8), Inches(2.0), Inches(3.8), Inches(1.5))
        self._add_metric_box(slide, "Total Videos", youtube_data.get("uploads", "N/A"), Inches(9.0), Inches(2.0), Inches(3.8), Inches(1.5))
        
        # Add additional information
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(2.5))
        info_text = info_box.text_frame
        
        p = info_text.add_paragraph()
        p.text = "Additional Information"
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(18)
        
        p = info_text.add_paragraph()
        p.text = f"Category: {youtube_data.get('category', 'N/A')}"
        
        p = info_text.add_paragraph()
        p.text = f"Channel Created: {youtube_data.get('created', 'N/A')}"
        
        p = info_text.add_paragraph()
        p.text = f"Estimated Monthly Earnings: {youtube_data.get('estimated_monthly_earnings', 'N/A')}"
        
        return slide
    
    def _add_tiktok_slide(self, tiktok_data):
        """Add TikTok statistics slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "TikTok Performance"
        
        # Clear default content placeholder
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx != 0:  # Skip title placeholder
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Add metric boxes
        self._add_metric_box(slide, "Followers", tiktok_data.get("followers", "N/A"), Inches(0.5), Inches(2.0), Inches(3.8), Inches(1.5))
        self._add_metric_box(slide, "Total Likes", tiktok_data.get("total_likes", "N/A"), Inches(4.8), Inches(2.0), Inches(3.8), Inches(1.5))
        self._add_metric_box(slide, "Videos", tiktok_data.get("videos", "N/A"), Inches(9.0), Inches(2.0), Inches(3.8), Inches(1.5))
        
        # Add additional information
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(2.5))
        info_text = info_box.text_frame
        
        p = info_text.add_paragraph()
        p.text = "Additional Information"
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(18)
        
        p = info_text.add_paragraph()
        p.text = f"Engagement Rate: {tiktok_data.get('engagement_rate', 'N/A')}"
        
        return slide
    
    def _add_suggestions_slides(self, suggestions):
        """Add optimization suggestions slides to the presentation"""
        # Create first suggestions slide
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Optimization Suggestions"
        
        # Add suggestions as bullet points
        content_shape = slide.placeholders[1]
        content_text = content_shape.text_frame
        
        # Add up to 5 suggestions to the first slide
        first_slide_suggestions = suggestions[:5] if len(suggestions) > 5 else suggestions
        
        for i, suggestion in enumerate(first_slide_suggestions):
            if i == 0:
                content_text.text = suggestion
            else:
                p = content_text.add_paragraph()
                p.text = suggestion
        
        # If we have more than 5 suggestions, create an additional slide
        if len(suggestions) > 5:
            additional_slide = self.presentation.slides.add_slide(slide_layout)
            additional_title = additional_slide.shapes.title
            additional_title.text = "Additional Recommendations"
            
            additional_content = additional_slide.placeholders[1]
            additional_text = additional_content.text_frame
            
            for i, suggestion in enumerate(suggestions[5:]):
                if i == 0:
                    additional_text.text = suggestion
                else:
                    p = additional_text.add_paragraph()
                    p.text = suggestion
        
        return slide
    
    def _add_next_steps_slide(self):
        """Add next steps slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Next Steps"
        
        # Set content
        content_shape = slide.placeholders[1]
        content_text = content_shape.text_frame
        
        steps = [
            "Review these insights and recommendations",
            "Create an action plan based on the suggestions",
            "Implement changes to your content strategy",
            "Monitor performance metrics",
            "Reassess in 30-60 days"
        ]
        
        for i, step in enumerate(steps):
            if i == 0:
                content_text.text = step
            else:
                p = content_text.add_paragraph()
                p.text = step
        
        return slide
    
    def _add_thank_you_slide(self):
        """Add thank you slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Thank You"
        
        # Set subtitle
        current_date = datetime.now().strftime("%B %d, %Y")
        subtitle = slide.placeholders[1]
        subtitle.text = f"Report generated on {current_date}"
        
        # Add contact information
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.0))
        contact_text = contact_box.text_frame
        contact_text.text = "For questions or additional analysis, please contact us."
        
        return slide
    
    def _add_metric_box(self, slide, title, value, left, top, width, height):
        """Add a metric box to the slide"""
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Add title text
        title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.5))
        title_text = title_box.text_frame
        title_text.text = title
        title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_text.paragraphs[0].runs[0].font.size = Pt(14)
        title_text.paragraphs[0].runs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        # Add value text
        value_box = slide.shapes.add_textbox(left, top + Inches(0.5), width, Inches(0.8))
        value_text = value_box.text_frame
        value_text.text = value
        value_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_text.paragraphs[0].runs[0].font.size = Pt(28)
        value_text.paragraphs[0].runs[0].font.bold = True
        value_text.paragraphs[0].runs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    def create_content_creator_presentation(self, analysis_results, creator_name=None):
        """
        Create a presentation based on content creator analysis results
        
        Args:
            analysis_results (dict): Results from the content creator analysis
            creator_name (str, optional): Name of the content creator
        
        Returns:
            bytes or dict: Presentation file as bytes or error message
        """
        # Generate a title for the presentation
        if creator_name:
            title = f"Content Strategy Analysis for {creator_name}"
        else:
            platforms = []
            if "youtube" in analysis_results and analysis_results["youtube"] and "error" not in analysis_results["youtube"]:
                platforms.append("YouTube")
            if "tiktok" in analysis_results and analysis_results["tiktok"] and "error" not in analysis_results["tiktok"]:
                platforms.append("TikTok")
            
            if platforms:
                title = f"Content Strategy Analysis: {' & '.join(platforms)}"
            else:
                title = "Content Creator Performance Analysis"
        
        # Generate the presentation
        presentation = self.generate_presentation(
            title=title,
            content=analysis_results,
            creator_name=creator_name
        )
        
        return presentation
    
    def get_download_link(self, presentation_bytes, filename="presentation.pptx"):
        """
        Generate a download link for the presentation
        
        Args:
            presentation_bytes (bytes): PowerPoint presentation file as bytes
            filename (str): Name of the file to download
            
        Returns:
            str: HTML download link
        """
        b64 = base64.b64encode(presentation_bytes).decode()
        href = f'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}'
        return href